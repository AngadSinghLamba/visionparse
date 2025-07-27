"""
📄 file_loader.py

This module handles loading and parsing of various file types:
- PDF: converts to images
- PNG/JPG/JPEG: loads image
- DOCX: extracts paragraphs
- PPTX: extracts slide text
- CSV/XLS: reads as text
"""

from pathlib import Path
from pdf2image import convert_from_path
from PIL import Image
import pandas as pd
import docx
from pptx import Presentation

# 🧠 Main function to handle all supported file types
def load_file(file_path):
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return load_pdf(file_path)
    elif suffix in [".png", ".jpg", ".jpeg"]:
        return load_image(file_path)
    elif suffix == ".docx":
        return load_docx(file_path)
    elif suffix == ".pptx":
        return load_pptx(file_path)
    elif suffix in [".csv", ".xls", ".xlsx"]:
        return load_spreadsheet(file_path)
    else:
        return f"❌ Unsupported file type: {suffix}"

# 📄 Load PDF and return list of images
def load_pdf(file_path):
    try:
        images = convert_from_path(str(file_path), dpi=300)
        return images
    except Exception as e:
        return f"❌ Error reading PDF: {e}"

# 🖼️ Load image file
def load_image(file_path):
    try:
        return [Image.open(file_path).convert("RGB")]
    except Exception as e:
        return f"❌ Error reading image: {e}"

# 📝 Load DOCX and extract paragraph text
def load_docx(file_path):
    try:
        doc = docx.Document(file_path)
        paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        return "\n".join(paragraphs)
    except Exception as e:
        return f"❌ Error reading DOCX: {e}"

# 📊 Load PPTX and extract text from all slides
def load_pptx(file_path):
    try:
        prs = Presentation(file_path)
        slide_texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_texts.append(shape.text.strip())
        return "\n".join(slide_texts)
    except Exception as e:
        return f"❌ Error reading PPTX: {e}"

# 📈 Load CSV or Excel and convert to markdown
def load_spreadsheet(file_path):
    try:
        if file_path.suffix == ".csv":
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
        return df.to_markdown(index=False)
    except Exception as e:
        return f"❌ Error reading spreadsheet: {e}"
